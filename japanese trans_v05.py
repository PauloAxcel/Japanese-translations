# -*- coding: utf-8 -*-
"""
Created on Tue Jul 25 09:03:47 2023

@author: Paulo Gomes
"""
import cv2
import pytesseract
from PIL import ImageGrab
from googletrans import Translator
import numpy as np
import mss
import pykakasi
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import numpy as np
from io import BytesIO
import os
from pptx.dml.color import RGBColor
from tqdm import tqdm
import MeCab
import pandas as pd
from googletrans import Translator
import win32api
x, y = win32api.GetCursorPos()

import nltk
nltk.download('words')
from nltk.corpus import words
import japanize_matplotlib
word_list = set(words.words())
import re

def acquire_screenshot(monitor_number):
    with mss.mss() as sct:
        mon = sct.monitors[monitor_number]
        screenshot = sct.grab(mon)
        screenshot = np.array(screenshot)
        # screenshot2 = screenshot[screen_region[0]:screen_region[1], screen_region[2]:screen_region[3]]
        print('\n Image acquired!... Quick change screen! ')
    
    return screenshot

def preprocess_image(image):
    # Apply image preprocessing techniques like resizing, noise reduction, and thresholding
    # Use OpenCV for these tasks
    # For example, convert the image to grayscale and apply thresholding
    gray_image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    _, threshold_image = cv2.threshold(gray_image, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
    return threshold_image

def perform_ocr(image):
    # Perform OCR using pytesseract on the preprocessed image
    # Set the correct language for Japanese characters
    # extracted_text = pytesseract.image_to_string(image, lang='jpn')
    extracted_text_data = pytesseract.image_to_data(image, lang='jpn', output_type=pytesseract.Output.DICT)
    return extracted_text_data

def translate_text(cleaned_extracted_text_data):
    # Use a translation API to translate the Japanese text
    # For example, using the googletrans library:
    translator = Translator()
    cleaned_extracted_text_data_trans = cleaned_extracted_text_data
    for cnt, matrix in enumerate(cleaned_extracted_text_data):
        text = matrix[0]
        translated_text = translator.translate(text, dest='en').text
        cleaned_extracted_text_data_trans[cnt].append(translated_text)
        
    print('\n translate_text completed!')
        
    return cleaned_extracted_text_data_trans



def clean_japanese_text(extracted_text_data):
    
    text = ('\n').join(extracted_text_data['text'])
    japanese_pattern = re.compile(r'[^\u3000-\u303F\u3040-\u309F\u30A0-\u30FF\uFF66-\uFF9F\u4E00-\u9FFF\n]+')
    
    # Use the regex pattern to remove non-Japanese characters from the text
    cleaned_text = japanese_pattern.sub('', text)
    
    # Remove the single consecutive newline character
    cleaned_text = re.sub(r'(?<=\S)\n(?=\S)', '', cleaned_text)

    # Replace multiple consecutive newline characters with a single newline character
    cleaned_text = re.sub(r'\n{2,}', '\n', cleaned_text)

    return cleaned_text

def convert_space_inside(lst):
    updated_list = []
    for element in lst:
        if element == ' ' or element == '  ':
            updated_list.append('')
        else:
            updated_list.append(element)
    return updated_list

def separate_sentences(text1,left,top,width,height):
    sentences = []
    tops = []
    lefts = []
    widths = []
    heights = []
    
    current_sentence = ""
    current_top = []
    current_left = []
    current_width = []
    current_height = []
    
    sentence_count = 1
    for cnt, (te,l,t,w,h) in enumerate(zip(text1,left,top,width,height)):
        if te != "":
            # Belongs to a sentence
            current_sentence += te + " "
            current_top.append(t)
            current_left.append(l)
            current_width.append(w)
            current_height.append(h)
        else:
            # Empty string, sentence ended
            if current_sentence.strip() != "":
                sentences.append(current_sentence.strip())
                tops.append([int(np.min(current_top)),int(np.max(current_top))])
                lefts.append([int(np.min(current_left)),int(np.max(current_left))])
                widths.append([int(np.min(current_width)),int(np.max(current_width))])
                heights.append([int(np.min(current_height)),int(np.max(current_height))])
                                        

                
                sentence_count += 1
                
                current_sentence = ""
                current_top = []
                current_left = []
                current_width = []
                current_height = []
    # Add the last sentence if it exists
    if current_sentence.strip() != "":
        sentences.append(current_sentence.strip())
        tops.append([int(np.min(current_top)),int(np.max(current_top))])
        lefts.append([int(np.min(current_left)),int(np.max(current_left))])
        widths.append([int(np.min(current_width)),int(np.max(current_width))])
        heights.append([int(np.min(current_height)),int(np.max(current_height))])
        
    return sentences, tops, lefts, widths, heights


def clean_extracted_data(extracted_text_data, cleaned_text, image_np):
    
    text0 = extracted_text_data['text']
    text1 = convert_space_inside(text0)
    
    
    left = extracted_text_data['left']
    top = extracted_text_data['top']
    width = extracted_text_data['width']
    height = extracted_text_data['height']
    
    sentences, tops, lefts, widths, heights = separate_sentences(text1,left,top,width,height)

    japanese_pattern = re.compile(r'[^\u3000-\u303F\u3040-\u309F\u30A0-\u30FF\uFF66-\uFF9F\u4E00-\u9FFF0-9\n]+')
    
    cleaned_text = [japanese_pattern.sub('', a) for a in sentences]
    
    spaces2 = [a != '' for a in cleaned_text]
    
    left2 = [a for s,a in zip(spaces2,lefts) if s == True]
    top2 = [a for s,a in zip(spaces2,tops) if s == True]
    width2 = [a for s,a in zip(spaces2,widths) if s == True]
    height2 = [a for s,a in zip(spaces2,heights) if s == True]
    text2 = [a for s,a in zip(spaces2,cleaned_text) if s == True]
    
    matrix = []
    for l,t,w,h,te in zip(left2,top2,width2,height2,text2):
        # Crop the corresponding region from the image
        # try:
        #     # Calculate the coordinates of the cropping region
        delta = 15
        crop_left = max(0, l[0] - int(np.min(w)) - delta)
        crop_top = max(0, t[0] + int(np.min(h)) - delta)
        crop_right = min(image_np.shape[1] - 1, l[1] + int(np.min(w)) + delta)
        crop_bottom = min(image_np.shape[0] - 1, t[1] + int(np.min(h)) + delta)
        matrix.append([te,crop_top ,crop_bottom, crop_left ,crop_right ])

        

    return matrix


def kanji_to_hiragana(text):
    # Use Google Translate for converting Kanji to Hiragana
    kks = pykakasi.kakasi()
    result = kks.convert(text)
    orig = []
    hira = []
    kata = []
    for item in result:
        orig.append(item['orig'])
        hira.append(item['hira'])
        kata.append(item['kana'])
        
    hiragana_text = ('').join(hira)

    return hiragana_text




from googletrans import Translator, LANGUAGES
import MeCab

def analyze_and_translate(text):
    mecab_tagger = MeCab.Tagger('-Owakati')
    tagger = MeCab.Tagger()
    translator = Translator()

    parsed_text = mecab_tagger.parse(text)
    parsed = parsed_text.strip().split()

    data = []
    max_pos_components = 0

    for cnt, word in enumerate(parsed):
        result = tagger.parse(word)
        lines = result.strip().split('\n')

        for line in lines:
            if line == 'EOS':
                break
            surface, reading, base, _, pos, _, _, position, *_ = line.split('\t')

            try:
                trans_surface = translator.translate(surface, dest='en').text
            except Exception as e:
                print(f"Translation error for '{surface}': {e}")
                trans_surface = "Translation not available"

            pos_components = pos.split('-')
            pos_components_trans = []
            for a in pos_components:
                try:
                    a_trans = translator.translate(a, dest='en').text
                except Exception as e:
                    print(f"Translation error for '{a}': {e}")
                    a_trans = "Translation not available"
                pos_components_trans.append(a_trans)

            data.append([str(cnt+1), surface, reading, base, trans_surface] + pos_components_trans)
            max_pos_components = max(max_pos_components, len(pos_components_trans))

    columns = ['Pos', 'S Form', 'Read', 'B Form', 'Trans'] + [f'Mean_{i+1}' for i in range(max_pos_components)]
    df = pd.DataFrame(data, columns=columns)

    return df




def organizer(image_np,cleaned_extracted_text_data_trans):
    # Display the original Japanese text and the translated text
    images = []
    texts = []
    dffs = []
    for matrix in tqdm(cleaned_extracted_text_data_trans):
        japanese_text = matrix[0]
        result_df = analyze_and_translate(japanese_text)
        english_text = matrix[-1]
        hiragana_text = kanji_to_hiragana(japanese_text)

        crop_top = matrix[1]
        crop_bottom = matrix[2]
        crop_left = matrix[3]
        crop_right = matrix[4]
        
        images.append(image_np[crop_top:crop_bottom, crop_left:crop_right])
        texts.append(japanese_text + '\n' + hiragana_text + '\n' + english_text)
        dffs.append(result_df)
    
    print('\n organized completed!')
        
    return images, texts, dffs
    



def save_results_to_ppt(images, texts, dffs, output_file):
    if os.path.exists(output_file):
        # If the file exists, open it for editing
        prs = Presentation(output_file)
    else:
        # If the file does not exist, create a new presentation
        prs = Presentation()
    
    for i, image in enumerate(images):
        # Get the corresponding title
        title = texts[i]
        
        # Add a slide to the presentation
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title to the slide
        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        # Change title font size to 28
        for paragraph in title_placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(28)

        # Define the position and size of the image on the slide
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)  # Change image width to 9 inches
        height = Inches(0.7)  # Change image height to 0.7 inches
        
        image_stream = BytesIO()
        Image.fromarray(image).save(image_stream, format='PNG')
        image_stream.seek(0)
        # Add the image to the slide
        slide.shapes.add_picture(image_stream, left, top, width, height)

        # Get the corresponding result_df for the image
        result_df = dffs[i]

        # Add a table to the slide
        table_left = Inches(1.5)
        table_top = Inches(5.5)
        table_width = Inches(8)
        table_height = Inches(0.8)
        table = slide.shapes.add_table(rows=result_df.shape[0]+1, cols=result_df.shape[1], 
                                       left=table_left, top=table_top,
                                       width=table_width, height=table_height).table

        # Set table column widths
        for j in range(result_df.shape[1]):
            table.columns[j].width = Inches(1.5)

        # Set table header format and font size
        for j, header_text in enumerate(result_df.columns):
            cell = table.cell(0, j)
            cell.text = header_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Set header font color to black
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if len(cell.text_frame.paragraphs[0].runs) == 0:
                cell.text_frame.paragraphs[0].add_run()
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
            cell.text_frame.word_wrap = True

        # Set table data format and font size for specific columns
        for row in range(result_df.shape[0]):
            for col in range(result_df.shape[1]):
                cell = table.cell(row+1, col)
                cell.text = str(result_df.iloc[row, col])
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.TOP
                if len(cell.text_frame.paragraphs[0].runs) == 0:
                    cell.text_frame.paragraphs[0].add_run()
                
                # Set font size based on the column index
                if col in [1, 2, 3]:  # Columns 2, 3, and 4 will have font size 16
                    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
                else:
                    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

                cell.text_frame.word_wrap = True

    # Save the PowerPoint presentation
    prs.save(output_file)



def main():   
    # screen_region = (ytop, ybottom, xtop, xbottom)  # Example region (top-left corner: 0,0, bottom-right corner: 800,600)
    monitor_number = 2
    pytesseract.pytesseract.tesseract_cmd = r'C:\Users\Paulo Gomes\AppData\Local\Programs\Tesseract-OCR\tesseract.exe'
    help_region = [700,900]
    # Capture the screen
    screen = acquire_screenshot(monitor_number)
    
    screen2 = screen[help_region[0]:help_region[1],:]
       
    # Preprocess the image
    processed_image = preprocess_image(screen2)
    
    # Perform OCR
    extracted_text_data = perform_ocr(processed_image)

    # Clean the text and keep only Japanese characters
    cleaned_text = clean_japanese_text(extracted_text_data)
    
    # [te,crop_top ,crop_bottom, crop_left ,crop_right ]
    cleaned_extracted_text_data = clean_extracted_data(extracted_text_data, cleaned_text, screen2)

    # Translate the cleaned text
    cleaned_extracted_text_data_trans = translate_text(cleaned_extracted_text_data)
    
    images, texts, dffs = organizer(screen2,cleaned_extracted_text_data_trans)
    
    output_file = 'japanese_output.pptx'
    save_results_to_ppt(images, texts, dffs, output_file) 

    return screen

if __name__ == "__main__":
    screen = main()



    
    
    
 






    
